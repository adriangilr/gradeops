from __future__ import annotations

import csv
import json
import io
import os
import copy
import re
import shutil
import zipfile
import unicodedata
from datetime import datetime, timedelta, timezone
from typing import Any

from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload

from src.auth import get_credentials
from src.config import ensure_directories, get_settings

from src.utils.naming import construir_nombre_portfolio as build_portfolio_name
from src.config import NAMING_MODE, MAX_FOLDER_NAME_LEN


try:
    from PyPDF2 import PdfReader  # type: ignore
except Exception:
    PdfReader = None

try:
    from docx import Document  # type: ignore
except Exception:
    Document = None

try:
    from pptx import Presentation  # type: ignore
except Exception:
    Presentation = None

# ==========================================================
# Configuración general
# ==========================================================

DEBUG_RUBRICS = False
RECENT_DAYS = 30
AUTOGRADING_CONFIG_FILENAME = "autograding_rules.json"
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tiff", ".tif"}

DEFAULT_AUTOGRADING_CONFIG = {
    "weights": {
        "delivery_valid": 40,
        "evidence_file_or_text": 20,
        "readable_content": 20,
        "minimum_sufficiency": 20,
    },
    "keywords": {
        "enabled": True,
        "list": ["control", "sistema", "resumen"],
        "minimum_matches": 1,
        "required_for_delivery_valid": False,
    },
    "minimum_sufficiency": {
        "min_words_partial": 10,
        "min_words_full": 50,
        "min_chars_partial": 80,
        "min_chars_full": 300,
        "partial_score": 10,
        "full_score": 20,
    },
    "late_policy": {
        "enabled": True,
        "minor_days_threshold": 5,
        "minor_penalty": 5,
        "major_penalty": 10,
        "fallback_penalty_when_late_without_due_date": 5,
    },
}

# Mapeo simple de MIME -> extensión para corregir nombres raros
MIME_TO_EXT = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/zip": ".zip",
    "text/plain": ".txt",
    "text/csv": ".csv",
    "application/json": ".json",
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/jpg": ".jpg",
    "image/webp": ".webp",
    "image/gif": ".gif",
    "image/tiff": ".tiff",
}




DEFAULT_LANGUAGE = "es"

def get_runtime_language() -> str:
    """
    Futuro punto central para cargar idioma desde config/runtime.
    """
    return DEFAULT_LANGUAGE


LANG = get_runtime_language()

I18N = {
    "es": {
        "ui": {
            "select": "Selecciona {type_label}:",
            "enter_number": "Ingresa número:",
            "go_back": "Regresar",
            "invalid_option": "❌ Opción inválida. Intenta de new.",
            "final_confirmation": "CONFIRMACION FINAL",
            "course": "Curso",
            "scope": "Alcance",
            "activity_filter": "Filtro de actividades",
            "download_mode": "Modo de descarga",
            "output_format": "Formato de salida",
            "activity": "Actividad",
            "activities_to_process": "Actividades a procesar",
            "confirm_download": "¿Confirmar descarga? (s/n) [s]: ",
            "download_cancelled": "⏹️ Descarga cancelada por el usuario.",
            "invalid_yes_no": "❌ Respuesta inválida. Escribe 's' o 'n'.",
        },
        "menu": {
            "download_scope": "el alcance de la descarga",
            "download_mode": "un modo de descarga",
            "output_format": "un formato de salida",
            "activity_filter": "un filter_name de actividades",
            "course": "un curso",
            "activity": "una actividad",
        },
        "options": {
            "single_coursework": "Descargar una sola actividad",
            "all_courseworks": "Descargar todas las actividades del curso",
            "all": "Incluir todos los alumnos (entregados y no entregados)",
            "late_ungraded": "Bajar solo tardías y no evaluadas",
            "folder_only": "Guardar solo en carpeta",
            "zip_and_folder": "Guardar en carpeta y generar .zip",
            "with_submissions": "Solo actividades con entregas",
            "all_activities": "Todas las actividades",
        },
        "descriptions": {
            "all": "todas las entregas activas",
            "resubmitted": "solo reentregadas y pendientes de reevaluación",
            "ungraded": "solo no evaluadas",
            "late": "solo tardías",
            "resubmitted_ungraded": "solo reentregadas y no evaluadas",
            "late_ungraded": "solo tardías y no evaluadas",
            "unknown_filter": "filter_name desconocido",
        },
        "runtime": {
            "config_read_error": "⚠️ No se pudo leer la configuración de autograding en '{path}': {err}",
            "downloaded": "      ✅ Descargado: {name}",
            "download_error": "      ❌ Error al descargar '{name}': {err}",
            "attachments_none": "  attachments: ninguno",
            "attachments_header": "  attachments:",
            "link": "    - Link: {title} | url={url}",
            "form": "    - Form: {title} | url={url}",
            "youtube": "    - YouTube: {title} | url={url}",
            "unhandled_attachment": "    - Tipo de adjunto no manejado directamente.",
            "csv_generated": "✅ CSV generado: {path}",
            "zip_generated": "✅ ZIP generado: {path}",
            "processing_activity": "Procesando actividad: {name}",
            "no_submissions": "No se encontraron entregas para esta actividad.",
            "submissions_total": "Entregas totales encontradas: {count}",
            "submissions_filtered": "Entregas a procesar con este filter_name: {count}",
            "submissions_no_match": "No hay entregas que coincidan con el filter_name para esta actividad.",
            "user_profile_fallback": "  ⚠️ No se pudo leer userProfiles para userId={user_id}. Se usará fallback. Detalle: {err}",
            "no_submission_download": "  attachments: no aplica, alumno sin entrega enviada",
            "auth_ok": "Autenticación correcta.",
            "token_valid": "Token válido: {value}",
            "profile_scope_ok": "✅ Scope de perfiles disponible.",
            "profile_scope_missing": "⚠️ El token no tiene scope para leer perfiles de alumnos. Se continuará con fallback y el CSV puede traer name/correo vacíos.",
            "no_courses": "No se encontraron cursos.",
            "selected_course": "✅ Curso seleccionado: {value}",
            "selected_scope": "📚 Alcance seleccionado: {value}",
            "selected_filter": "🧩 Filtro de actividades: {value}",
            "selected_mode": "📥 Modo seleccionado: {value}",
            "selected_output": "📦 Formato de salida: {value}",
            "no_courseworks": "No se encontraron actividades en este curso.",
            "no_courseworks_after_filter": "No quedaron actividades después de aplicar el filter_name.",
            "courseworks_found": "✅ Actividades encontradas: {total} | después del filter_name: {filtered}",
            "selected_activity": "✅ Actividad seleccionada: {value}",
            "processing_all_courseworks": "✅ Se procesarán todas las actividades filtered_items del curso: {count}",
            "unknown_scope": "❌ Alcance de descarga no reconocido.",
            "final_summary": "RESUMEN FINAL",
            "activities_processed": "Actividades procesadas: {count}",
            "submissions_seen": "Entregas totales vistas: {count}",
            "submissions_filter_matched": "Entregas que cumplieron filter_name: {count}",
            "files_downloaded": "Archivos descargados: {count}",
            "csv_rows": "Filas en CSV: {count}",
            "all_mode_note": "Nota: en modo 'all' ahora el CSV incluye también alumnos sin entregar.",
            "base_folder": "Carpeta base: {path}",
            "zip_path": "ZIP: {path}",
            "process_done": "✅ Proceso terminado.",
            "classroom_error": "Error al consultar Classroom: {err}",
            "deliverable_id": "Entrega ID: {value}",
            "user_id": "  userId: {value}",
            "state": "  estado: {value}",
            "late": "  late: {value}",
            "assigned_grade": "  assignedGrade: {value}",
            "draft_grade": "  draftGrade: {value}",
            "resubmitted": "  reentregada: {value}",
            "ungraded": "  no evaluada: {value}",
            "attached": "  attached: {value}",
        },
        "feedback": {
            "no_submission": "Sin entrega enviada. Calificacion automatica en 0.",
            "manual_review": "Requiere revision manual por type_label de evidencia.",
            "readable_content": "Contenido legible detectado con {word_count} palabras aprox.",
            "not_readable": "No se detecto contenido legible automaticamente.",
            "keyword_hits": "Coincidencias clave: {keywords}.",
            "late_submission": "Entrega tardia de {days_late} dia(s).",
            "suggested_grade": "Calificacion automatica sugerida: {auto_score}/100.",
            "reason_no_submission": "No submitted. Score final 0.",
            "valid_delivery": "Entrega valida con evidencia detectada",
            "no_evidence": "Sin evidencia adjunta",
            "readable": "contenido legible",
            "not_readable_short": "sin contenido legible",
            "sufficiency_points": "suficiencia {points} pts",
            "sufficiency_zero": "suficiencia 0 pts",
            "keywords_detected": "keywords detectadas: {keywords}",
            "late_penalty": "penalizacion por tardanza {points} pts",
            "manual_review_marked": "marcada para revision manual",
            "evidence_received": "Se recibió evidencia de entrega.",
            "no_interpretable_evidence": "No se detectó evidencia adjunta ni text interpretable.",
            "late_penalty_applied": "Se aplicó penalización por tardanza (-{points}).",
            "no_late_penalty": "Sin penalización por tardanza.",
            "manual_review_long": "La evidencia requiere revision manual porque incluye imagenes o formatos no interpretables automaticamente.",
            "auto_readable_detected": "Se detectó contenido legible automáticamente ({word_count} palabras aprox.).",
            "evidence_not_interpretable": "Se recibio evidencia, pero no fue posible interpretarla automaticamente con las librerias actuales.",
            "no_readable_for_auto": "No fue posible recuperar contenido legible para evaluation automatica.",
            "keywords_detected_long": "Palabras clave detectadas: {keywords}.",
            "full_sufficiency": "El contenido cumple suficiencia completa según los umbrales configurados.",
            "partial_sufficiency": "El contenido cumple suficiencia mínima parcial; conviene revisión rápida.",
            "low_sufficiency": "El contenido detectado es breve o insuficiente según los umbrales configurados.",
            "compact_late": "tardía",
            "compact_manual_review": "manual_review",
            "compact_not_readable": "no_legible",
            "compact_full": "suficiente",
            "compact_partial": "parcial",
            "compact_low": "insuficiente",
            "compact_valid": "valida",
        },
        "fallbacks": {
            "unnamed_course": "Curso sin name",
            "unnamed_activity": "Actividad sin título",
            "untitled": "sin_titulo",
            "course": "curso",
            "activity": "actividad",
            "file": "archivo",
            "no_user_id": "sin_userId",
            "no_id": "Sin ID",
            "no_state": "Sin estado",
            "no_url": "Sin URL",
            "no_title": "Sin título",
            "room": "Aula",
        },
        "labels": {
            "submission_status": {
                "TURNED_IN": "turned_in",
                "CREATED": "assigned_not_submitted",
                "RETURNED": "returned",
                "RECLAIMED_BY_STUDENT": "reclaimed_by_student",
                "UNKNOWN": "unknown",
            },
            "submission_type": {
                "none": "none",
                "text_only": "text_only",
                "image_only": "image_only",
                "file_only": "file_only",
                "mixed": "mixed",
                "link_only": "link_only",
            },
            "bool": {
                True: "true",
                False: "false",
            },
        },
    },
    "en": {
        "ui": {
            "select": "Select {type_label}:",
            "enter_number": "Enter number:",
            "go_back": "Go back",
            "invalid_option": "❌ Invalid option. Try again.",
            "final_confirmation": "FINAL CONFIRMATION",
            "course": "Course",
            "scope": "Scope",
            "activity_filter": "Activity filter",
            "download_mode": "Download mode",
            "output_format": "Output format",
            "activity": "Activity",
            "activities_to_process": "Activities to process",
            "confirm_download": "Confirm download? (y/n) [y]: ",
            "download_cancelled": "⏹️ Download cancelled by user.",
            "invalid_yes_no": "❌ Invalid response. Type 'y' or 'n'.",
        },
        "menu": {
            "download_scope": "the download scope",
            "download_mode": "a download mode",
            "output_format": "an output format",
            "activity_filter": "an activity filter",
            "course": "a course",
            "activity": "an activity",
        },
        "options": {
            "single_coursework": "Download a single activity",
            "all_courseworks": "Download all course activities",
            "all": "Include all students (submitted and not submitted)",
            "late_ungraded": "Download only late and ungraded",
            "folder_only": "Save folder only",
            "zip_and_folder": "Save folder and generate .zip",
            "with_submissions": "Only activities with submissions",
            "all_activities": "All activities",
        },
        "descriptions": {
            "all": "all active submissions",
            "resubmitted": "only resubmitted and pending reevaluation",
            "ungraded": "only ungraded",
            "late": "only late",
            "resubmitted_ungraded": "only resubmitted and ungraded",
            "late_ungraded": "only late and ungraded",
            "unknown_filter": "unknown filter",
        },
        "runtime": {
            "config_read_error": "⚠️ Could not read autograding config at '{path}': {err}",
            "downloaded": "      ✅ Downloaded: {name}",
            "download_error": "      ❌ Error downloading '{name}': {err}",
            "attachments_none": "  attachments: none",
            "attachments_header": "  attachments:",
            "link": "    - Link: {title} | url={url}",
            "form": "    - Form: {title} | url={url}",
            "youtube": "    - YouTube: {title} | url={url}",
            "unhandled_attachment": "    - Attachment type not directly handled.",
            "csv_generated": "✅ CSV generated: {path}",
            "zip_generated": "✅ ZIP generated: {path}",
            "processing_activity": "Processing activity: {name}",
            "no_submissions": "No submissions found for this activity.",
            "submissions_total": "Total submissions found: {count}",
            "submissions_filtered": "Submissions to process with this filter: {count}",
            "submissions_no_match": "No submissions match this filter for this activity.",
            "user_profile_fallback": "  ⚠️ Could not read userProfiles for userId={user_id}. Fallback will be used. Detail: {err}",
            "no_submission_download": "  attachments: not applicable, student did not submit",
            "auth_ok": "Authentication successful.",
            "token_valid": "Valid token: {value}",
            "profile_scope_ok": "✅ Profile scope available.",
            "profile_scope_missing": "⚠️ Token has no scope to read student profiles. Continuing with fallback and CSV may contain empty name/email fields.",
            "no_courses": "No courses found.",
            "selected_course": "✅ Selected course: {value}",
            "selected_scope": "📚 Selected scope: {value}",
            "selected_filter": "🧩 Activity filter: {value}",
            "selected_mode": "📥 Selected mode: {value}",
            "selected_output": "📦 Output format: {value}",
            "no_courseworks": "No activities found in this course.",
            "no_courseworks_after_filter": "No activities remained after applying the filter.",
            "courseworks_found": "✅ Activities found: {total} | after filter: {filtered}",
            "selected_activity": "✅ Selected activity: {value}",
            "processing_all_courseworks": "✅ All filtered course activities will be processed: {count}",
            "unknown_scope": "❌ Unknown download scope.",
            "final_summary": "FINAL SUMMARY",
            "activities_processed": "Activities processed: {count}",
            "submissions_seen": "Total submissions seen: {count}",
            "submissions_filter_matched": "Submissions matching filter: {count}",
            "files_downloaded": "Files downloaded: {count}",
            "csv_rows": "CSV rows: {count}",
            "all_mode_note": "Note: in 'all' mode the CSV now also includes students with no submission.",
            "base_folder": "Base folder: {path}",
            "zip_path": "ZIP: {path}",
            "process_done": "✅ Process finished.",
            "classroom_error": "Error querying Classroom: {err}",
            "deliverable_id": "Submission ID: {value}",
            "user_id": "  userId: {value}",
            "state": "  state: {value}",
            "late": "  late: {value}",
            "assigned_grade": "  assignedGrade: {value}",
            "draft_grade": "  draftGrade: {value}",
            "resubmitted": "  resubmitted: {value}",
            "ungraded": "  ungraded: {value}",
            "attached": "  attached: {value}",
        },
        "feedback": {
            "no_submission": "No submission sent. Automatic grade set to 0.",
            "manual_review": "Requires manual review due to evidence type.",
            "readable_content": "Readable content detected with about {word_count} words.",
            "not_readable": "No readable content was detected automatically.",
            "keyword_hits": "Keyword matches: {keywords}.",
            "late_submission": "Late submission by {days_late} day(s).",
            "suggested_grade": "Suggested automatic grade: {auto_score}/100.",
            "reason_no_submission": "No submission. Final score 0.",
            "valid_delivery": "Valid submission with detected evidence",
            "no_evidence": "No attached evidence",
            "readable": "readable content",
            "not_readable_short": "no readable content",
            "sufficiency_points": "sufficiency {points} pts",
            "sufficiency_zero": "sufficiency 0 pts",
            "keywords_detected": "keywords detected: {keywords}",
            "late_penalty": "late penalty {points} pts",
            "manual_review_marked": "marked for manual review",
            "evidence_received": "Submission evidence was received.",
            "no_interpretable_evidence": "No attached evidence or interpretable text was detected.",
            "late_penalty_applied": "Late penalty applied (-{points}).",
            "no_late_penalty": "No late penalty.",
            "manual_review_long": "Evidence requires manual review because it includes images or formats that cannot be interpreted automatically.",
            "auto_readable_detected": "Readable content detected automatically ({word_count} words approx.).",
            "evidence_not_interpretable": "Evidence was received, but could not be interpreted automatically with the current libraries.",
            "no_readable_for_auto": "Could not recover readable content for automatic evaluation.",
            "keywords_detected_long": "Keywords detected: {keywords}.",
            "full_sufficiency": "Content meets full sufficiency according to configured thresholds.",
            "partial_sufficiency": "Content meets partial minimum sufficiency; a quick review is recommended.",
            "low_sufficiency": "Detected content is brief or insufficient according to configured thresholds.",
            "compact_late": "late",
            "compact_manual_review": "manual_review",
            "compact_not_readable": "not_readable",
            "compact_full": "sufficient",
            "compact_partial": "partial",
            "compact_low": "insufficient",
            "compact_valid": "valid",
        },
        "fallbacks": {
            "unnamed_course": "Untitled course",
            "unnamed_activity": "Untitled activity",
            "untitled": "untitled",
            "course": "course",
            "activity": "activity",
            "file": "file",
            "no_user_id": "no_userId",
            "no_id": "No ID",
            "no_state": "No state",
            "no_url": "No URL",
            "no_title": "Untitled",
            "room": "Room",
        },
        "labels": {
            "submission_status": {
                "TURNED_IN": "turned_in",
                "CREATED": "assigned_not_submitted",
                "RETURNED": "returned",
                "RECLAIMED_BY_STUDENT": "reclaimed_by_student",
                "UNKNOWN": "unknown",
            },
            "submission_type": {
                "none": "none",
                "text_only": "text_only",
                "image_only": "image_only",
                "file_only": "file_only",
                "mixed": "mixed",
                "link_only": "link_only",
            },
            "bool": {
                True: "true",
                False: "false",
            },
        },
    },
}


class _SafeFormatDict(dict):
    def __missing__(self, key: str) -> str:
        return "{" + key + "}"


def t(key: str, lang: str | None = None, **kwargs: Any) -> str:
    active_lang = lang or LANG
    fallback_lang = "es"
    try:
        value: Any = I18N[active_lang]
        for part in key.split("."):
            value = value[part]
    except Exception:
        value = I18N[fallback_lang]
        for part in key.split("."):
            value = value[part]

    if not isinstance(value, str):
        return str(value)

    format_kwargs = dict(kwargs)
    if "type_label" in format_kwargs and "tipo" not in format_kwargs:
        format_kwargs["tipo"] = format_kwargs["type_label"]
    if "name" in format_kwargs and "nombre" not in format_kwargs:
        format_kwargs["nombre"] = format_kwargs["name"]
    if "last_name" in format_kwargs and "apellido" not in format_kwargs:
        format_kwargs["apellido"] = format_kwargs["last_name"]

    return value.format_map(_SafeFormatDict(format_kwargs)) if format_kwargs else value


def labels() -> dict[str, Any]:
    return I18N.get(LANG, I18N["es"])["labels"]



# ==========================================================
# CSV schema e idiomas configurables
# ==========================================================

CSV_OUTPUT_COLUMNS = [
    "course_name",
    "activity_name",
    "student_name",
    "student_mail",
    "submission_status",
    "has_attachment",
    "submission_type",
    "primary_file_type",
    "days_late",
    "is_readable",
    "word_count",
    "keyword_hits",
    "requires_manual_review",
    "confidence_score",
    "auto_score",
    "auto_feedback",
    "auto_grading_reason",
    "ai_feedback",
    "final_grade",
    "final_feedback",
]



def normalize_basic_ascii(text: str) -> str:
    """
    Deja text simple, sin acentos ni caracteres especiales problematicos.
    """
    text = unicodedata.normalize("NFKD", text or "")
    text = text.encode("ascii", "ignore").decode("ascii")
    text = re.sub(r"\s+", " ", text).strip()
    return text


def bool_to_text(value: bool) -> str:
    return labels()["bool"][bool(value)]


def get_submission_status(submission: dict[str, Any]) -> str:
    state = (submission.get("state") or "").upper()
    return labels()["submission_status"].get(
        state,
        labels()["submission_status"]["UNKNOWN"],
    )


def detect_submission_type(
    submission: dict[str, Any],
    downloaded_paths: list[str],
    readable_content: bool,
) -> str:
    """
    Clasifica el type_label de entrega con una heuristica simple y mantenible.
    """
    attachments = get_attachments(submission)

    if not attachments and readable_content:
        return labels()["submission_type"]["text_only"]

    if not attachments and not downloaded_paths:
        return labels()["submission_type"]["none"]

    tiene_imagen = contains_images(downloaded_paths)
    tiene_archivo_no_imagen = any(
        os.path.splitext(path)[1].lower() not in IMAGE_EXTENSIONS
        for path in downloaded_paths
    )

    tiene_links = any(
        ("link" in att) or ("form" in att) or ("youTubeVideo" in att)
        for att in attachments
    )

    if tiene_imagen and tiene_archivo_no_imagen:
        return labels()["submission_type"]["mixed"]
    if tiene_imagen and not tiene_archivo_no_imagen:
        return labels()["submission_type"]["image_only"]
    if tiene_archivo_no_imagen:
        return labels()["submission_type"]["file_only"]
    if tiene_links:
        return labels()["submission_type"]["link_only"]
    return labels()["submission_type"]["none"]


def calculate_confidence_score(
    submitted: bool,
    has_attachment: bool,
    readable_content: bool,
    manual_review: bool,
    word_count: int,
    keyword_hits_count: int,
    submission_type: str,
) -> float:
    """
    Score heuristico 0..1 para indicar confianza del autograding.
    """
    if not submitted:
        return 1.0

    score = 0.35
    if has_attachment:
        score += 0.15
    if readable_content:
        score += 0.25
    if word_count >= 50:
        score += 0.15
    elif word_count >= 10:
        score += 0.08

    if keyword_hits_count > 0:
        score += min(0.10, keyword_hits_count * 0.05)

    if submission_type == labels()["submission_type"]["mixed"]:
        score -= 0.10
    if submission_type == labels()["submission_type"]["image_only"]:
        score -= 0.25
    if manual_review:
        score -= 0.30

    return round(max(0.0, min(1.0, score)), 2)


def build_auto_feedback(
    submitted: bool,
    readable_content: bool,
    manual_review: bool,
    word_count: int,
    days_late: int,
    keyword_hits: list[str],
    auto_score: int,
) -> str:
    if not submitted:
        return t("feedback.no_submission")

    parts: list[str] = []

    if manual_review:
        parts.append(t("feedback.manual_review"))
    elif readable_content:
        parts.append(t("feedback.readable_content", word_count=word_count))
    else:
        parts.append(t("feedback.not_readable"))

    if keyword_hits:
        parts.append(t("feedback.keyword_hits", keywords=", ".join(keyword_hits)))

    if days_late > 0:
        parts.append(t("feedback.late_submission", days_late=days_late))

    parts.append(t("feedback.suggested_grade", auto_score=auto_score))
    return normalize_basic_ascii(" ".join(parts))


def build_auto_grading_reason(
    submitted: bool,
    has_attachment: bool,
    readable_content: bool,
    sufficiency_score: int,
    keyword_hits: list[str],
    late_penalty: int,
    manual_review: bool,
) -> str:
    if not submitted:
        return t("feedback.reason_no_submission")

    reasons: list[str] = []

    if has_attachment:
        reasons.append(t("feedback.valid_delivery"))
    else:
        reasons.append(t("feedback.no_evidence"))

    if readable_content:
        reasons.append(t("feedback.readable"))
    else:
        reasons.append(t("feedback.not_readable_short"))

    if sufficiency_score > 0:
        reasons.append(t("feedback.sufficiency_points", points=sufficiency_score))
    else:
        reasons.append(t("feedback.sufficiency_zero"))

    if keyword_hits:
        reasons.append(t("feedback.keywords_detected", keywords=", ".join(keyword_hits)))

    if late_penalty > 0:
        reasons.append(t("feedback.late_penalty", points=late_penalty))

    if manual_review:
        reasons.append(t("feedback.manual_review_marked"))

    return normalize_basic_ascii("; ".join(reasons) + ".")


# ==========================================================
# Utilidades generales
# ==========================================================

def sanitize_file_name(text: str) -> str:
    """
    Limpia nombres de archivos o carpetas.
    Evita 'sin_nombre' usando fallback más inteligente.
    """
    if not text or not text.strip():
        from datetime import datetime
        return f"{t('fallbacks.file')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

    text = text.strip()
    replacements = {
        "/": "-",
        "\\": "-",
        ":": "-",
        "*": "-",
        "?": "",
        '"': "",
        "<": "(",
        ">": ")",
        "|": "-",
    }

    for old, new in replacements.items():
        text = text.replace(old, new)

    import re
    text = re.sub(r"\s+", " ", text).strip()

    if not text:
        from datetime import datetime
        return f"{t('fallbacks.file')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}"

    return text


def slugify_name(text: str) -> str:
    """
    Convierte text a formato carpeta amigable:
    minusculas, guiones y sin caracteres raros.
    """
    text = sanitize_file_name(text).lower()
    text = text.replace("_", " ")
    text = re.sub(r"[^a-z0-9]+", "-", text)
    text = re.sub(r"-+", "-", text).strip("-")
    return text or "sin-name"


def build_course_slug(text: str, course_id: str) -> str:
    """
    Genera el name visible de la carpeta raíz del curso.
    Ejemplo:
    'Seminario Ing. de SW' + id -> 'Seminario Ing. de SW_840182924161'
    """
    nombre_curso = sanitize_file_name(text) or t("fallbacks.course")
    return f"{nombre_curso}_{course_id}"


def build_activity_slug(text: str, actividad_id: str) -> str:
    """
    Genera el name visible de la carpeta de la actividad.
    Ejemplo:
    'P01 - Timing plan' + id -> 'P01 - Timing plan_840182924183'
    """
    nombre_actividad = sanitize_file_name(text) or t("fallbacks.activity")
    return f"{nombre_actividad}_{actividad_id}"


def get_folder_timestamp() -> str:
    """
    Timestamp corto y estable para nombres de carpeta.
    Formato: YYYYMMDD_HHMM
    """
    return datetime.now().strftime("%Y%m%d_%H%M")

def ensure_directory(path: str) -> None:
    """
    Crea el directorio si no existe.
    """
    os.makedirs(path, exist_ok=True)


def prepare_output_directory(path: str, limpiar_si_existe: bool = False) -> str:
    """
    Prepara un directorio de salida controlado.
    Si limpiar_si_existe=True, elimina por completo el contenido previo.
    """
    path = os.path.normpath(path)

    if limpiar_si_existe and os.path.exists(path):
        shutil.rmtree(path)

    ensure_directory(path)
    return path


def ensure_extension(name: str, mime_type: str) -> str:
    """
    Garantiza que el name tenga una extensión coherente con el MIME.
    Esto evita casos donde Classroom/Drive entrega títulos raros como .pod
    o nombres sin extensión.
    """
    base, ext = os.path.splitext(name)
    ext_actual = ext.lower()
    ext_correcta = MIME_TO_EXT.get((mime_type or "").lower())

    if ext_correcta is None:
        return name

    if not ext_actual:
        return f"{base}{ext_correcta}"

    if ext_actual != ext_correcta:
        return f"{base}{ext_correcta}"

    return name


def build_submission_folder_name(
    submission: dict,
    profile: dict,
) -> str:

    name = profile.get("name", "")
    last_name = profile.get("last_name", "")
    user_id = str(submission.get("userId", ""))

    return build_portfolio_name(
        nombre=name,
        apellido=last_name,
        user_id=user_id,
        modo=NAMING_MODE,
        max_len=MAX_FOLDER_NAME_LEN,
    )


def select_option(
    items: list[dict[str, Any]],
    type_label: str,
    allow_back: bool = False,
) -> dict[str, Any] | None:
    """
    Menú interactivo genérico para terminal.
    Si allow_back=True, muestra una opción extra para volver.
    """
    while True:
        print(f"\n{t('ui.select', tipo=type_label)}\n")

        for i, item in enumerate(items, start=1):
            print(f"{i}. {item['display_name']}")

        back_option = len(items) + 1
        if allow_back:
            print(f"{back_option}. {t('ui.go_back')}")

        input_value = input(f"\n{t('ui.enter_number')} ").strip()

        try:
            index = int(input_value) - 1

            if 0 <= index < len(items):
                return items[index]

            if allow_back and index == len(items):
                return None
        except ValueError:
            pass

        print(t("ui.invalid_option"))


def parse_google_datetime(value: str | None) -> datetime | None:
    """
    Convierte timestamps de Google type_label:
    2026-04-09T18:25:43.123Z
    """
    if not value:
        return None

    try:
        if value.endswith("Z"):
            value = value.replace("Z", "+00:00")
        return datetime.fromisoformat(value)
    except ValueError:
        return None


def utc_now() -> datetime:
    """
    Regresa la fecha actual en UTC.
    """
    return datetime.now(timezone.utc)


def load_autograding_config() -> dict[str, Any]:
    """
    Carga reglas de autograding desde JSON.
    Si no existe o falla, usa la configuración por defecto.
    """
    possible_paths = [
        os.path.join(os.path.dirname(__file__), AUTOGRADING_CONFIG_FILENAME),
        os.path.join(os.getcwd(), AUTOGRADING_CONFIG_FILENAME),
    ]

    for path in possible_paths:
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8") as f:
                    data = json.load(f)

                if isinstance(data, dict):
                    return merge_config(DEFAULT_AUTOGRADING_CONFIG, data)
            except Exception as err:
                print(t("runtime.config_read_error", path=path, err=err))

    return copy.deepcopy(DEFAULT_AUTOGRADING_CONFIG)


def merge_config(base: dict[str, Any], extra: dict[str, Any]) -> dict[str, Any]:
    """Mezcla profunda simple de diccionarios."""
    result: dict[str, Any] = {}

    for key, value in base.items():
        if isinstance(value, dict):
            extra_value = extra.get(key, {})
            if isinstance(extra_value, dict):
                result[key] = merge_config(value, extra_value)
            else:
                result[key] = copy.deepcopy(value)
        else:
            result[key] = extra.get(key, value)

    for key, value in extra.items():
        if key not in result:
            result[key] = value

    return result



def get_rubric_schema_version(workbook) -> str:
    """
    Lee schema_version desde tab CONFIG del Rubric.xlsx.

    Busca:
    Column B -> "Schema Version"
    Column C -> valor

    Fallback:
    rubric_runtime_v1
    """
    try:
        if "CONFIG" not in workbook.sheetnames:
            return "rubric_runtime_v1"

        sheet = workbook["CONFIG"]

        for row in sheet.iter_rows(values_only=True):
            if not row or len(row) < 3:
                continue

            key = str(row[1] or "").strip().lower()

            if key in {
                "schema version",
                "schema_version",
                "version schema",
            }:
                value = str(row[2] or "").strip()

                if value:
                    return value

    except Exception:
        pass

    return "rubric_runtime_v1"


def build_rubric_runtime_json(
    rubric_xlsx_path: str = "config/Rubric.xlsx",
    output_json_path: str = "config/rubric_runtime.json",
    debug: bool = True,
) -> dict[str, Any]:
    """
    Convierte Rubric.xlsx a un JSON interno estandarizado.

    Mejora clave:
    - Detecta automaticamente la fila real de headers.
    - Normaliza headers: mayusculas, espacios, acentos y variantes comunes.
    - No rompe el flujo actual: solo genera config/rubric_runtime.json.
    """

    try:
        from openpyxl import load_workbook
    except Exception as err:
        raise RuntimeError(
            "openpyxl no instalado. Ejecuta: python -m pip install openpyxl"
        ) from err

    if not os.path.exists(rubric_xlsx_path):
        print(f"⚠️ Rubric XLSX no encontrado: {rubric_xlsx_path}")
        return {
            "schema_version": schema_version,
            "generated_at": datetime.now().isoformat(),
            "source_file": rubric_xlsx_path,
            "rubrics": [],
        }

    def normalize_header(value: Any) -> str:
        raw = normalize_basic_ascii(str(value or "").strip().lower())
        raw = raw.replace("-", "_").replace("/", "_").replace(".", "_")
        raw = re.sub(r"[^a-z0-9_]+", "_", raw)
        raw = re.sub(r"_+", "_", raw).strip("_")

        aliases = {
            "categoria": "category",
            "category": "category",

            "criterio": "criterion_name",
            "nombre_criterio": "criterion_name",
            "criterion": "criterion_name",
            "criterion_name": "criterion_name",

            "id_criterio": "criterion_id",
            "criterio_id": "criterion_id",
            "criterion_id": "criterion_id",

            "id_rubrica": "rubric_id",
            "rubrica_id": "rubric_id",
            "rubric_id": "rubric_id",

            "puntaje_obtenido": "obtained_score",
            "calificacion_obtenida": "obtained_score",
            "obtained_score": "obtained_score",

            "puntaje_maximo": "max_score",
            "calificacion_maxima": "max_score",
            "max_score": "max_score",

            "keywords": "matched_keywords",
            "keyword_hits": "matched_keywords",
            "matched_keywords": "matched_keywords",
            "palabras_clave": "matched_keywords",

            "estado": "status",
            "status": "status",

            "revision_manual": "manual_review",
            "manual_review": "manual_review",
        }

        level_aliases = {
            "4_accomplished": "4-accomplished",
            "4_accomplished_": "4-accomplished",
            "accomplished": "4-accomplished",
            "logrado": "4-accomplished",

            "3_competent": "3-competent",
            "competent": "3-competent",
            "competente": "3-competent",

            "2_developing": "2-developing",
            "developing": "2-developing",
            "en_desarrollo": "2-developing",

            "1_beginning": "1-beginning",
            "beginning": "1-beginning",
            "inicial": "1-beginning",

            "0_not_accomplished": "0-not_accomplished",
            "not_accomplished": "0-not_accomplished",
            "no_logrado": "0-not_accomplished",
        }

        return aliases.get(raw, level_aliases.get(raw, raw))

    def normalize_headers(row: tuple[Any, ...]) -> list[str]:
        headers: list[str] = []
        used: dict[str, int] = {}

        for cell in row:
            header = normalize_header(cell)
            if not header:
                headers.append("")
                continue

            # Evita headers duplicados sin romper el parseo.
            if header in used:
                used[header] += 1
                header = f"{header}_{used[header]}"
            else:
                used[header] = 1

            headers.append(header)

        return headers

    def find_header_row(rows: list[tuple[Any, ...]]) -> tuple[int, list[str]]:
        """
        Busca la fila de headers en las primeras filas.
        Esto soporta tabs con titulo, metadata o instrucciones antes de la tabla.
        """
        expected = {
            "category",
            "criterion_name",
            "criterion_id",
            "rubric_id",
            "max_score",
            "matched_keywords",
        }

        best_index = 0
        best_headers: list[str] = []
        best_score = -1

        for idx, row in enumerate(rows[:25]):
            headers = normalize_headers(row)
            score = len(set(headers) & expected)

            # Header minimo util: criterion_name o criterion_id.
            if score > best_score:
                best_index = idx
                best_headers = headers
                best_score = score

        return best_index, best_headers

    workbook = load_workbook(rubric_xlsx_path, data_only=True)

    schema_version = get_rubric_schema_version(workbook)

    runtime_data: dict[str, Any] = {
        "schema_version": "rubric_runtime_v1",
        "generated_at": datetime.now().isoformat(),
        "source_file": rubric_xlsx_path,
        "rubrics": [],
    }

    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))

        if not rows:
            continue

        header_row_index, headers = find_header_row(rows)        
        rubric = {
            "sheet_name": sheet.title,
            "rubric_id": None,
            "criteria": [],
        }

        for row in rows[header_row_index + 1:]:
            if row is None:
                continue

            item = {
                headers[i]: row[i]
                for i in range(min(len(headers), len(row)))
                if headers[i]
            }

            criterion_id = str(item.get("criterion_id") or "").strip()
            criterion_name = str(item.get("criterion_name") or "").strip()
            category = str(item.get("category") or "").strip()

            # Ignora filas vacias, notas o instrucciones sin criterio.
            if not criterion_id and not criterion_name and not category:
                continue

            # Evita que tabs como CONFIG/README se vuelvan criterios basura.
            if not criterion_id and not criterion_name:
                continue

            rubric_id = str(item.get("rubric_id") or "").strip()
            if not rubric_id:
                rubric_id = sheet.title if sheet.title.upper() not in {"CONFIG", "README"} else ""

            if rubric["rubric_id"] is None and rubric_id:
                rubric["rubric_id"] = rubric_id

            criterion_runtime = {
                "criterion_id": criterion_id,
                "rubric_id": rubric_id,
                "category": category,
                "criterion_name": criterion_name,
                "max_score": safe_int(item.get("max_score"), 0),

                "levels": {
                    "4": str(item.get("4-accomplished") or "").strip(),
                    "3": str(item.get("3-competent") or "").strip(),
                    "2": str(item.get("2-developing") or "").strip(),
                    "1": str(item.get("1-beginning") or "").strip(),
                    "0": str(item.get("0-not_accomplished") or "").strip(),
                },

                "matched_keywords": parse_keywords_runtime(
                    item.get("matched_keywords")
                ),
                "manual_review": normalize_bool(
                    item.get("manual_review")
                ),
                "obtained_score": safe_int(
                    item.get("obtained_score"),
                    0,
                ),
                "status": str(item.get("status") or "").strip(),

                "compatibility_mode": "hybrid",
                "language": LANG,
            }

            rubric["criteria"].append(criterion_runtime)
        
        criteria_detected = len(rubric["criteria"])
        
        if DEBUG_RUBRICS:
            print(f"\n📄 Sheet: {sheet.title}")
            print(f"   header_row: {header_row_index + 1}")
            print(f"   headers: {headers}")
            print(f"   criteria_detected: {criteria_detected}")

        if rubric["criteria"]:
            runtime_data["rubrics"].append(rubric)


    output_dir = os.path.dirname(output_json_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)

    with open(output_json_path, "w", encoding="utf-8") as f:
        json.dump(
            runtime_data,
            f,
            indent=4,
            ensure_ascii=False,
        )

    print(f"✅ Rubric runtime JSON generado: {output_json_path}")

    return runtime_data

def safe_int(value: Any, default: int = 0) -> int:
    try:
        if value is None or value == "":
            return default
        return int(float(value))
    except Exception:
        return default


def normalize_bool(value: Any) -> bool:
    if value is None:
        return False

    text = str(value).strip().lower()

    return text in {
        "true",
        "yes",
        "y",
        "1",
        "si",
        "sí",
    }


def parse_keywords_runtime(value: Any) -> list[dict[str, Any]]:
    """
    Convierte:
    diseno-2,sistema-2/modulo-3

    a:

    [
        {
            "type": "AND",
            "keywords": [
                {"term": "diseno", "min_occurrences": 2},
                {"term": "sistema", "min_occurrences": 2},
            ]
        },
        {
            "type": "OR",
            "keywords": [
                {"term": "modulo", "min_occurrences": 3},
            ]
        }
    ]
    """

    if value is None:
        return []

    text = str(value).strip()

    if not text or text.upper() == "N/A":
        return []

    groups = []

    or_groups = text.split("/")

    for group in or_groups:

        group = group.strip()

        if not group:
            continue

        and_keywords = []

        for token in group.split(","):

            token = token.strip()

            if not token:
                continue

            if "-" in token:
                term, count = token.rsplit("-", 1)

                and_keywords.append({
                    "term": normalize_basic_ascii(term.strip().lower()),
                    "min_occurrences": safe_int(count, 1),
                })

            else:
                and_keywords.append({
                    "term": normalize_basic_ascii(token.lower()),
                    "min_occurrences": 1,
                })

        groups.append({
            "type": "AND_GROUP",
            "keywords": and_keywords,
        })

    return groups




















# ==========================================================
# Nombres visibles para menús
# ==========================================================

def get_visible_course_name(course: dict[str, Any]) -> str:
    """
    Construye un name visible único para evitar confusiones
    cuando existen dos cursos con el mismo name.
    """
    name = course.get("name", t("fallbacks.unnamed_course"))
    section = course.get("section", "").strip()
    course_id = course.get("id", "").strip()
    room = course.get("room", "").strip()

    extras = []
    if section:
        extras.append(section)
    if room:
        extras.append(f"{t('fallbacks.room')}: {room}")
    if course_id:
        extras.append(f"id={course_id}")

    return f"{name} | {' | '.join(extras)}" if extras else name


def get_visible_activity_name(coursework: dict[str, Any]) -> str:
    """
    Construye un name más legible para la actividad.
    """
    title = coursework.get("title", t("fallbacks.unnamed_activity"))
    coursework_id = coursework.get("id", "")
    work_type = coursework.get("workType", "")
    max_points = coursework.get("maxPoints")

    extras = []
    if work_type:
        extras.append(work_type)
    if max_points is not None:
        extras.append(f"{max_points} pts")
    if coursework_id:
        extras.append(f"id={coursework_id}")

    return f"{title} | {' | '.join(extras)}" if extras else title


# ==========================================================
# Menús
# ==========================================================

def select_download_scope(allow_back: bool = False) -> str | None:
    """
    Permite elegir si se descargará una actividad
    o todas las actividades del curso.
    """
    opciones = [
        {
            "id": "single_coursework",
            "display_name": t("options.single_coursework"),
        },
        {
            "id": "all_courseworks",
            "display_name": t("options.all_courseworks"),
        },
    ]
    seleccion = select_option(
        opciones,
        t("menu.download_scope"),
        allow_back=allow_back,
    )
    return seleccion["id"] if seleccion else None


def select_download_mode(allow_back: bool = False) -> str | None:
    """
    Define qué entregas se descargarán.
    """
    opciones = [
        {"id": "all", "display_name": t("options.all")},
        {
            "id": "late_ungraded",
            "display_name": t("options.late_ungraded"),
        },
    ]
    seleccion = select_option(
        opciones,
        t("menu.download_mode"),
        allow_back=allow_back,
    )
    return seleccion["id"] if seleccion else None


def select_output_format(allow_back: bool = False) -> str | None:
    """
    Define si solo se guarda en carpeta o también en zip.
    """
    opciones = [
        {"id": "folder_only", "display_name": t("options.folder_only")},
        {"id": "zip_and_folder", "display_name": t("options.zip_and_folder")},
    ]
    seleccion = select_option(
        opciones,
        t("menu.output_format"),
        allow_back=allow_back,
    )
    return seleccion["id"] if seleccion else None


def select_activity_filter(allow_back: bool = False) -> str | None:
    """
    Permite aplicar un filter_name a las actividades del curso
    antes de elegir una o antes de procesarlas todas.
    """
    opciones = [
        {"id": "all", "display_name": t("options.all_activities")},
        {"id": "with_submissions", "display_name": t("options.with_submissions")},
    ]
    seleccion = select_option(
        opciones,
        t("menu.activity_filter"),
        allow_back=allow_back,
    )
    return seleccion["id"] if seleccion else None


def describe_download_mode(download_mode: str) -> str:
    descripciones = {
        "all": t("descriptions.all"),
        "resubmitted": t("descriptions.resubmitted"),
        "ungraded": t("descriptions.ungraded"),
        "late": t("descriptions.late"),
        "resubmitted_ungraded": t("descriptions.resubmitted_ungraded"),
        "late_ungraded": t("descriptions.late_ungraded"),
    }
    return descripciones.get(download_mode, t("descriptions.unknown_filter"))


def confirm_download_summary(
    course_display: str,
    download_scope: str,
    activity_filter: str,
    download_mode: str,
    output_format: str,
    total_actividades: int,
    selected_activity: dict[str, Any] | None = None,
) -> bool:
    """
    Pide confirmación final antes de iniciar la descarga real.
    """
    print("\n" + "=" * 90)
    print(t("ui.final_confirmation"))
    print("=" * 90)
    print(f"{t('ui.course')}: {course_display}")
    print(f"{t('ui.scope')}: {download_scope}")
    print(f"{t('ui.activity_filter')}: {activity_filter}")
    print(f"{t('ui.download_mode')}: {describe_download_mode(download_mode)}")
    print(f"{t('ui.output_format')}: {output_format}")

    if selected_activity is not None:
        print(f"{t('ui.activity')}: {selected_activity.get('display_name', selected_activity.get('title', t('fallbacks.untitled')))}")
    else:
        print(f"{t('ui.activities_to_process')}: {total_actividades}")

    
        
        
        
    while True:
        input_value = input(f"\n{t('ui.confirm_download')}").strip().lower()

        if input_value == "":
            return True
        if input_value in {"s", "si", "sí", "y", "yes"}:
            return True
        if input_value in {"n", "no"}:
            print(t("ui.download_cancelled"))
            return False

        print(t("ui.invalid_yes_no"))


def get_export_directory(settings) -> str:
    """
    Directorio único para salidas finales del proceso.
    Centraliza todas las exportaciones en out/
    """
    return os.path.normpath("out")


# ==========================================================
# Descarga de Drive
# ==========================================================

def get_drive_file_metadata(drive_service, file_id: str) -> dict[str, str]:
    """
    Lee metadata real del archivo en Drive para usar el name correcto
    y el mimeType real, en lugar de confiar ciegamente en 'title'.
    """
    try:
        meta = (
            drive_service.files()
            .get(fileId=file_id, fields="id,name,mimeType,fileExtension")
            .execute()
        )
        return {
            "name": meta.get("name", ""),
            "mimeType": meta.get("mimeType", ""),
            "fileExtension": meta.get("fileExtension", ""),
        }
    except HttpError:
        return {
            "name": "",
            "mimeType": "",
            "fileExtension": "",
        }


def download_file(
    drive_service,
    file_id: str,
    file_name: str,
    folder: str,
    mime_type: str = "",
) -> str | None:
    """
    Descarga un archivo de Drive al folder indicado.
    Usa metadata real de Drive para corregir extensiones raras o faltantes.
    """
    ensure_directory(folder)

    try:
        meta = get_drive_file_metadata(drive_service, file_id)

        # Fuente de verdad:
        # 1) name real de Drive
        # 2) si no viene, usar file_name recibido
        real_name = meta.get("name") or file_name or t("fallbacks.file")
        real_mime = meta.get("mimeType") or mime_type or ""

        safe_name = sanitize_file_name(real_name)
        safe_name = ensure_extension(safe_name, real_mime)

        file_path = os.path.join(folder, safe_name)

        request = drive_service.files().get_media(fileId=file_id)
        with io.FileIO(file_path, "wb") as fh:
            downloader = MediaIoBaseDownload(fh, request)

            done = False
            while not done:
                _, done = downloader.next_chunk()

        print(t("runtime.downloaded", name=safe_name))
        return file_path

    except HttpError as err:
        print(t("runtime.download_error", name=file_name, err=err))
        return None


# ==========================================================
# Lectura de Classroom
# ==========================================================

def get_all_courses(classroom_service) -> list[dict[str, Any]]:
    """
    Recupera todos los cursos paginando.
    Si existen cursos activos, prioriza esos.
    """
    courses: list[dict[str, Any]] = []
    page_token = None

    while True:
        response = (
            classroom_service.courses()
            .list(pageSize=100, pageToken=page_token)
            .execute()
        )

        courses.extend(response.get("courses", []))
        page_token = response.get("nextPageToken")

        if not page_token:
            break

    cursos_activos = [c for c in courses if c.get("courseState", "ACTIVE") == "ACTIVE"]
    cursos_finales = cursos_activos if cursos_activos else courses

    for course in cursos_finales:
        course["display_name"] = get_visible_course_name(course)

    cursos_finales.sort(key=lambda x: x.get("display_name", "").lower())
    return cursos_finales


def get_all_activities(classroom_service, course_id: str) -> list[dict[str, Any]]:
    """
    Recupera todas las actividades de un curso.
    """
    courseworks: list[dict[str, Any]] = []
    page_token = None

    while True:
        response = (
            classroom_service.courses()
            .courseWork()
            .list(courseId=course_id, pageSize=100, pageToken=page_token)
            .execute()
        )

        courseworks.extend(response.get("courseWork", []))
        page_token = response.get("nextPageToken")

        if not page_token:
            break

    for coursework in courseworks:
        coursework["display_name"] = get_visible_activity_name(coursework)

    courseworks.sort(key=lambda x: x.get("display_name", "").lower())
    return courseworks


def get_all_submissions(
    classroom_service,
    course_id: str,
    coursework_id: str,
) -> list[dict[str, Any]]:
    """
    Recupera todas las entregas de una actividad.
    """
    submissions: list[dict[str, Any]] = []
    page_token = None

    while True:
        response = (
            classroom_service.courses()
            .courseWork()
            .studentSubmissions()
            .list(
                courseId=course_id,
                courseWorkId=coursework_id,
                pageSize=100,
                pageToken=page_token,
            )
            .execute()
        )

        submissions.extend(response.get("studentSubmissions", []))
        page_token = response.get("nextPageToken")

        if not page_token:
            break

    return submissions


# ==========================================================
# Filtros de entregas
# ==========================================================

def was_previously_returned(submission: dict[str, Any]) -> bool:
    """
    Detecta si la entrega ya había sido devuelta antes.
    """
    for event in submission.get("submissionHistory", []):
        state_history = event.get("stateHistory", {})
        if state_history.get("state") == "RETURNED":
            return True
    return False


def is_resubmitted(submission: dict[str, Any]) -> bool:
    return (
        submission.get("state") == "TURNED_IN"
        and was_previously_returned(submission)
    )


def is_ungraded(submission: dict[str, Any]) -> bool:
    return (
        submission.get("state") == "TURNED_IN"
        and submission.get("assignedGrade") is None
    )


def is_late(submission: dict[str, Any]) -> bool:
    return (
        submission.get("state") == "TURNED_IN"
        and submission.get("late", False) is True
    )


def get_readable_submission_status(submission: dict[str, Any]) -> str:
    """
    Traduce el estado tecnico de Classroom a un value configurable.
    """
    return get_submission_status(submission)


def can_download_submission(submission: dict[str, Any]) -> bool:
    """
    Solo tiene sentido intentar descargar attachments cuando la entrega fue enviada.
    """
    return submission.get("state") == "TURNED_IN"


def filter_submissions(submissions: list[dict[str, Any]], download_mode: str) -> list[dict[str, Any]]:
    """
    Aplica el filter_name elegido a la items de entregas.
    """
    if download_mode == "all":
        # Incluye a todos los alumnos de la actividad:
        # entregados, asignados sin entregar, devueltos, etc.
        return submissions

    if download_mode == "resubmitted":
        return [s for s in submissions if is_resubmitted(s)]

    if download_mode == "ungraded":
        return [s for s in submissions if is_ungraded(s)]

    if download_mode == "late":
        return [s for s in submissions if is_late(s)]

    if download_mode == "resubmitted_ungraded":
        return [s for s in submissions if is_resubmitted(s) and is_ungraded(s)]

    if download_mode == "late_ungraded":
        return [s for s in submissions if is_late(s) and is_ungraded(s)]

    return []


# ==========================================================
# Filtros de actividades
# ==========================================================

def is_published_activity(coursework: dict[str, Any]) -> bool:
    """
    Considera publicada cuando state es PUBLISHED.
    Si el campo no viene, asumimos True para no esconder actividades válidas.
    """
    state = coursework.get("state")
    if state is None:
        return True
    return state == "PUBLISHED"


def is_recent_activity(coursework: dict[str, Any], dias: int = RECENT_DAYS) -> bool:
    """
    Considera reciente una actividad creada o actualizada dentro
    de los últimos N días.
    """
    limit_dt = utc_now() - timedelta(days=dias)

    creation_time = parse_google_datetime(coursework.get("creationTime"))
    update_time = parse_google_datetime(coursework.get("updateTime"))
    due_date = None

    due = coursework.get("dueDate")
    if isinstance(due, dict):
        try:
            year = due.get("year")
            month = due.get("month")
            day = due.get("day")
            if year and month and day:
                due_date = datetime(year, month, day, tzinfo=timezone.utc)
        except ValueError:
            due_date = None

    valid_dates = [f for f in [creation_time, update_time, due_date] if f is not None]
    if not valid_dates:
        return False

    return any(f >= limit_dt for f in valid_dates)


def filter_activities(
    classroom_service,
    course_id: str,
    courseworks: list[dict[str, Any]],
    filter_name: str,
) -> list[dict[str, Any]]:
    """
    Aplica filter_name a las actividades.
    """
    if filter_name == "all":
        return courseworks

    if filter_name == "published":
        return [cw for cw in courseworks if is_published_activity(cw)]

    if filter_name == "recent":
        return [cw for cw in courseworks if is_recent_activity(cw)]

    if filter_name == "with_submissions":
        filtered_items = []
        for cw in courseworks:
            try:
                submissions = get_all_submissions(
                    classroom_service=classroom_service,
                    course_id=course_id,
                    coursework_id=cw["id"],
                )
                if submissions:
                    filtered_items.append(cw)
            except HttpError as err:
                print(
                    f"⚠️ No se pudieron revisar entregas para actividad "
                    f"'{cw.get('title', cw.get('id', 'sin_titulo'))}': {err}"
                )
        return filtered_items

    return courseworks


# ==========================================================
# Perfil del alumno
# ==========================================================

def extract_user_data_from_history(submission: dict[str, Any]) -> dict[str, str]:
    """
    Intenta recuperar name/correo desde submissionHistory.
    Esto sirve como plan B cuando no hay scopes suficientes
    para consultar userProfiles.
    """
    history = submission.get("submissionHistory", [])

    for event in history:
        actor = event.get("actorUser", {})
        if not actor:
            continue

        name = actor.get("name", {}) or {}
        given_name = name.get("givenName", "") or ""
        family_name = name.get("familyName", "") or ""
        full_name = name.get("fullName", "") or ""

        email = actor.get("emailAddress", "") or ""

        if not given_name and full_name:
            parts = full_name.split()
            if parts:
                given_name = parts[0]
                if len(parts) > 1:
                    family_name = " ".join(parts[1:])

        if given_name or family_name or email:
            return {
                "correo": email,
                "name": given_name,
                "last_name": family_name,
            }

    return {
        "correo": "",
        "name": "",
        "last_name": "",
    }


def get_user_profile(
    classroom_service,
    user_id: str,
    profile_scope_available: bool,
) -> dict[str, str]:
    """
    Recupera correo, name y last_name del alumno.

    Estrategia:
    1. Si hay scope disponible, intenta userProfiles
    2. Si no hay scope, el caller puede usar fallback desde historial
    """
    if not profile_scope_available:
        return {
            "correo": "",
            "name": "",
            "last_name": "",
        }

    try:
        profile = classroom_service.userProfiles().get(userId=user_id).execute()

        email = profile.get("emailAddress", "") or ""

        name = profile.get("name", {}) or {}
        given_name = name.get("givenName", "") or ""
        family_name = name.get("familyName", "") or ""
        full_name = name.get("fullName", "") or ""

        if not given_name and full_name:
            parts = full_name.split()
            if parts:
                given_name = parts[0]
                if len(parts) > 1:
                    family_name = " ".join(parts[1:])

        return {
            "correo": email,
            "name": given_name,
            "last_name": family_name,
        }

    except HttpError as err:
        raise err


def detect_profile_scope(classroom_service) -> bool:
    """
    Prueba una sola vez si el token tiene permisos para consultar perfiles.
    Así evitamos un error 403 repetido para cada alumno.
    """
    try:
        classroom_service.userProfiles().get(userId="me").execute()
        return True
    except HttpError as err:
        status = getattr(err, "status_code", None)
        contenido = str(err)

        if status == 403 or "ACCESS_TOKEN_SCOPE_INSUFFICIENT" in contenido:
            return False

        raise err


# ==========================================================
# Adjuntos
# ==========================================================

def get_attachments(submission: dict[str, Any]) -> list[dict[str, Any]]:
    assignment_submission = submission.get("assignmentSubmission", {})
    return assignment_submission.get("attachments", [])


def has_attachments(submission: dict[str, Any]) -> bool:
    return len(get_attachments(submission)) > 0


def get_due_date_text(coursework: dict[str, Any]) -> str:
    """
    Convierte dueDate de Classroom a text YYYY-MM-DD.
    """
    due_date = coursework.get("dueDate")
    if not isinstance(due_date, dict):
        return ""

    year = due_date.get("year")
    month = due_date.get("month")
    day = due_date.get("day")

    if not (year and month and day):
        return ""

    try:
        return f"{int(year):04d}-{int(month):02d}-{int(day):02d}"
    except (TypeError, ValueError):
        return ""


def get_due_time_text(coursework: dict[str, Any]) -> str:
    """
    Convierte dueTime de Classroom a text HH:MM:SS.
    """
    due_time = coursework.get("dueTime")
    if not isinstance(due_time, dict):
        return ""

    hours = due_time.get("hours", 0)
    minutes = due_time.get("minutes", 0)
    seconds = due_time.get("seconds", 0)

    try:
        return f"{int(hours):02d}:{int(minutes):02d}:{int(seconds):02d}"
    except (TypeError, ValueError):
        return ""


# ==========================================================
# Lectura y evaluación de contenido
# ==========================================================

def read_txt_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def read_pdf_text(path: str) -> str:
    if PdfReader is None:
        return ""

    try:
        reader = PdfReader(path)
        parts: list[str] = []
        for page in reader.pages:
            parts.append(page.extract_text() or "")
        return "\n".join(parts)
    except Exception:
        return ""


def read_docx_text(path: str) -> str:
    if Document is None:
        return ""

    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        return ""


def read_zip_text(path: str, profundidad_max: int = 15) -> str:
    """
    Intenta leer text útil de archivos simples dentro de un ZIP.
    No revienta si el ZIP trae binarios.
    """
    parts: list[str] = []

    try:
        with zipfile.ZipFile(path, "r") as zf:
            for idx, name in enumerate(zf.namelist()):
                if idx >= profundidad_max:
                    break

                lower = name.lower()
                if lower.endswith((".txt", ".md", ".csv", ".py", ".json", ".log")):
                    try:
                        data = zf.read(name)
                        parts.append(data.decode("utf-8", errors="ignore"))
                    except Exception:
                        continue
    except Exception:
        return ""

    return "\n".join(parts)


def is_image_file(path: str) -> bool:
    """
    Detecta si el archivo es una imagen común.
    """
    ext = os.path.splitext(path)[1].lower()
    return ext in IMAGE_EXTENSIONS


def contains_images(rutas: list[str]) -> bool:
    """
    Indica si entre los attachments descargados hay al menos una imagen.
    """
    return any(is_image_file(path) for path in rutas)


def get_primary_file_type(rutas: list[str]) -> str:
    """
    Regresa la extensión principal detectada en los archivos descargados.
    Si no hay archivos, regresa vacío.
    """
    if not rutas:
        return ""

    for path in rutas:
        ext = os.path.splitext(path)[1].lower().lstrip(".")
        if ext:
            return ext

    return ""


def read_pptx_text(path: str) -> str:
    if Presentation is None:
        return ""

    try:
        prs = Presentation(path)
        textos = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    textos.append(shape.text)

        return "\n".join(textos)
    except Exception:
        return ""


def extract_file_text(path: str) -> str:
    """
    Lee text de varios tipos de archivo comunes.
    Si no puede, regresa cadena vacía.
    """
    ext = os.path.splitext(path)[1].lower()

    if ext in {".txt", ".md", ".csv", ".py", ".json", ".log"}:
        return read_txt_text(path)

    if ext == ".pdf":
        return read_pdf_text(path)

    if ext == ".docx":
        return read_docx_text(path)

    if ext == ".zip":
        return read_zip_text(path)

    if ext == ".pptx":
        return read_pptx_text(path)

    if ext in IMAGE_EXTENSIONS:
        return ""

    return ""


def analyze_text_content(
    text: str,
    autograding_config: dict[str, Any],
) -> dict[str, Any]:
    """
    Analiza el text extraído con una lógica más simple:
    - detecta si el contenido es legible
    - calcula suficiencia mínima configurable
    - detecta palabras clave opcionales
    """
    clean_text = re.sub(r"\s+", " ", text or "").strip()
    palabras = re.findall(r"\b\w+\b", clean_text, flags=re.UNICODE)
    word_count = len(palabras)
    char_count = len(clean_text)

    weights = autograding_config.get("weights", {})
    suff_cfg = autograding_config.get("minimum_sufficiency", {})
    keywords_cfg = autograding_config.get("keywords", {})

    readable_content = bool(clean_text)

    min_words_partial = int(suff_cfg.get("min_words_partial", 10))
    min_words_full = int(suff_cfg.get("min_words_full", 50))
    min_chars_partial = int(suff_cfg.get("min_chars_partial", 80))
    min_chars_full = int(suff_cfg.get("min_chars_full", 300))
    partial_score = int(suff_cfg.get("partial_score", 10))
    full_score = int(suff_cfg.get("full_score", weights.get("minimum_sufficiency", 20)))

    if word_count >= min_words_full or char_count >= min_chars_full:
        sufficiency_score = full_score
        sufficiency_level = "full"
    elif word_count >= min_words_partial or char_count >= min_chars_partial:
        sufficiency_score = partial_score
        sufficiency_level = "partial"
    else:
        sufficiency_score = 0
        sufficiency_level = "low"

    keyword_list = keywords_cfg.get("list", [])
    if not isinstance(keyword_list, list):
        keyword_list = []

    keyword_hits: list[str] = []
    if keywords_cfg.get("enabled", True):
        texto_lower = clean_text.lower()
        for keyword in keyword_list:
            kw = str(keyword).strip().lower()
            if kw and kw in texto_lower:
                keyword_hits.append(kw)

    minimum_matches = int(keywords_cfg.get("minimum_matches", 1))
    keywords_ok = len(keyword_hits) >= minimum_matches if keyword_list else True

    return {
        "texto_extraido": clean_text,
        "word_count": word_count,
        "char_count": char_count,
        "readable_content": readable_content,
        "sufficiency_score": sufficiency_score,
        "sufficiency_level": sufficiency_level,
        "keyword_hits": keyword_hits,
        "keywords_ok": keywords_ok,
    }


def build_due_datetime(coursework: dict[str, Any]) -> datetime | None:
    due_date = coursework.get("dueDate")
    if not isinstance(due_date, dict):
        return None

    year = due_date.get("year")
    month = due_date.get("month")
    day = due_date.get("day")
    if not (year and month and day):
        return None

    due_time = coursework.get("dueTime") or {}
    hours = due_time.get("hours", 23)
    minutes = due_time.get("minutes", 59)
    seconds = due_time.get("seconds", 59)

    try:
        return datetime(
            int(year),
            int(month),
            int(day),
            int(hours),
            int(minutes),
            int(seconds),
            tzinfo=timezone.utc,
        )
    except (TypeError, ValueError):
        return None


def get_submission_timestamp(submission: dict[str, Any]) -> datetime | None:
    candidatos = [
        submission.get("updateTime"),
        submission.get("submissionTime"),
        submission.get("turnInTime"),
        submission.get("creationTime"),
    ]
    for value in candidatos:
        dt = parse_google_datetime(value)
        if dt is not None:
            return dt
    return None


def calculate_late_penalty(
    submission: dict[str, Any],
    coursework: dict[str, Any],
    autograding_config: dict[str, Any],
) -> tuple[int, int]:
    late_policy = autograding_config.get("late_policy", {})
    if not late_policy.get("enabled", True):
        return 0, 0

    if not bool(submission.get("late", False)):
        return 0, 0

    due_dt = build_due_datetime(coursework)
    submission_dt = get_submission_timestamp(submission)

    if due_dt is None or submission_dt is None:
        penalty = int(late_policy.get("fallback_penalty_when_late_without_due_date", 5))
        return penalty, 0

    delta = submission_dt - due_dt
    total_dias = max(0, int((delta.total_seconds() + 86399) // 86400))

    threshold = int(late_policy.get("minor_days_threshold", 5))
    minor_penalty = int(late_policy.get("minor_penalty", 5))
    major_penalty = int(late_policy.get("major_penalty", 10))

    penalty = minor_penalty if total_dias <= threshold else major_penalty
    return penalty, total_dias


def build_feedback(
    late: bool,
    has_attachment: bool,
    read_files: int,
    word_count: int,
    late_penalty: int,
    auto_grade: int,
    manual_review: bool,
    readable_content: bool,
    sufficiency_level: str,
    keyword_hits: list[str],
) -> str:
    """
    Genera feedback más limpio y profesional.
    Evita messages basura cuando sí existe evidencia, pero no fue interpretable.
    """
    messages: list[str] = []

    # 1) Evidencia
    if has_attachment:
        messages.append(t("feedback.evidence_received"))
    else:
        messages.append(t("feedback.no_interpretable_evidence"))

    # 2) Tardanza
    if late and late_penalty > 0:
        messages.append(t("feedback.late_penalty_applied", points=late_penalty))
    else:
        messages.append(t("feedback.no_late_penalty"))

    # 3) Estado de lectura
    if manual_review:
        messages.append(
            t("feedback.manual_review_long")
        )
    elif read_files > 0 and readable_content:
        messages.append(
            t("feedback.auto_readable_detected", num_palabras=word_count)
        )
    elif has_attachment:
        messages.append(
            t("feedback.evidence_not_interpretable")
        )
    else:
        messages.append(
            t("feedback.no_readable_for_auto")
        )

    # 4) Keywords: solo se informan si sí existen
    if keyword_hits:
        messages.append(t("feedback.keywords_detected_long", keywords=", ".join(keyword_hits)))

    # 5) Suficiencia
    if sufficiency_level == "full":
        messages.append(t("feedback.full_sufficiency"))
    elif sufficiency_level == "partial":
        messages.append(t("feedback.partial_sufficiency"))
    else:
        messages.append(t("feedback.low_sufficiency"))

    # 6) Cierre
    messages.append(t("feedback.suggested_grade", auto_score=auto_grade))
    return " ".join(messages)


def build_short_feedback(
    late: bool,
    manual_review: bool,
    readable_content: bool,
    sufficiency_level: str,
) -> str:
    """
    Genera feedback compacto type_label etiquetas para CSV y dashboards.
    Ejemplos:
    - suficiente
    - parcial
    - insuficiente
    - tardía | manual_review | insuficiente
    """
    tags: list[str] = []

    if late:
        tags.append(t("feedback.compact_late"))

    if manual_review:
        tags.append(t("feedback.compact_manual_review"))
    elif not readable_content:
        tags.append(t("feedback.compact_not_readable"))

    if sufficiency_level == "full":
        tags.append(t("feedback.compact_full"))
    elif sufficiency_level == "partial":
        tags.append(t("feedback.compact_partial"))
    else:
        tags.append(t("feedback.compact_low"))

    if not tags:
        return t("feedback.compact_valid")

    return " | ".join(tags)


def evaluate_submission_automatically(
    submission: dict[str, Any],
    downloaded_paths: list[str],
    coursework: dict[str, Any],
    autograding_config: dict[str, Any],
) -> dict[str, Any]:
    """
    Evalúa una entrega con estrategia simple y configurable.
    """
    weights = autograding_config.get("weights", {})
    keywords_cfg = autograding_config.get("keywords", {})

    late = bool(submission.get("late", False))
    submitted = submission.get("state") == "TURNED_IN"
    has_attachment = len(downloaded_paths) > 0 or has_attachments(submission)

    full_text_parts: list[str] = []
    read_files = 0
    manual_review = contains_images(downloaded_paths)

    for path in downloaded_paths:
        text = extract_file_text(path)
        if text.strip():
            full_text_parts.append(text)
            read_files += 1

    merged_text = "\n".join(full_text_parts)
    analysis = analyze_text_content(merged_text, autograding_config=autograding_config)

    late_penalty, days_late = calculate_late_penalty(
        submission=submission,
        coursework=coursework,
        autograding_config=autograding_config,
    )

    keyword_required_for_delivery = bool(
        keywords_cfg.get("required_for_delivery_valid", False)
    )

    delivery_valid_score = 0
    if submitted and (not keyword_required_for_delivery or analysis["keywords_ok"]):
        delivery_valid_score = int(weights.get("delivery_valid", 40))

    evidence_score = (
        int(weights.get("evidence_file_or_text", 20))
        if (has_attachment or analysis["readable_content"])
        else 0
    )
    readable_content_score = (
        int(weights.get("readable_content", 20))
        if analysis["readable_content"]
        else 0
    )
    sufficiency_score = int(analysis["sufficiency_score"])

    if not submitted:
        auto_grade = 0
        manual_review = False
    else:
        auto_grade = (
            delivery_valid_score
            + evidence_score
            + readable_content_score
            + sufficiency_score
            - late_penalty
        )
        auto_grade = max(0, min(100, auto_grade))

    keyword_hits_list = list(analysis["keyword_hits"])
    submission_type = detect_submission_type(
        submission=submission,
        downloaded_paths=downloaded_paths,
        readable_content=bool(analysis["readable_content"]),
    )
    primary_file_type = get_primary_file_type(downloaded_paths)
    confidence_score = calculate_confidence_score(
        submitted=submitted,
        has_attachment=has_attachment,
        readable_content=bool(analysis["readable_content"]),
        manual_review=manual_review,
        word_count=int(analysis["word_count"]),
        keyword_hits_count=len(keyword_hits_list),
        submission_type=submission_type,
    )
    auto_feedback = build_auto_feedback(
        submitted=submitted,
        readable_content=bool(analysis["readable_content"]),
        manual_review=manual_review,
        word_count=int(analysis["word_count"]),
        days_late=days_late,
        keyword_hits=keyword_hits_list,
        auto_score=auto_grade,
    )
    auto_grading_reason = build_auto_grading_reason(
        submitted=submitted,
        has_attachment=has_attachment,
        readable_content=bool(analysis["readable_content"]),
        sufficiency_score=sufficiency_score,
        keyword_hits=keyword_hits_list,
        late_penalty=late_penalty,
        manual_review=manual_review,
    )

    return {
        "auto_grade": auto_grade,
        "auto_score": auto_grade,
        "feedback": auto_feedback,
        "auto_feedback": auto_feedback,
        "auto_grading_reason": auto_grading_reason,
        "confidence_score": confidence_score,
        "submission_type": submission_type,
        "primary_file_type": primary_file_type,
        "late_penalty": late_penalty,
        "days_late": days_late,
        "has_attachment": bool_to_text(has_attachment),
        "delivery_valid_score": delivery_valid_score,
        "evidence_score": evidence_score,
        "readable_content_score": readable_content_score,
        "content_score": sufficiency_score,
        "minimum_sufficiency_score": sufficiency_score,
        "files_read_for_content": read_files,
        "detected_words": int(analysis["word_count"]),
        "detected_characters": int(analysis["char_count"]),
        "keyword_hits": ", ".join(keyword_hits_list),
        "keyword_hits_list": keyword_hits_list,
        "manual_review": bool_to_text(manual_review),
        "requires_manual_review": bool_to_text(manual_review),
        "readable_content": bool_to_text(bool(analysis["readable_content"])),
        "is_readable": bool_to_text(bool(analysis["readable_content"])),
    }


def print_submission_summary(submission: dict[str, Any]) -> None:
    """
    Imprime información útil para la revisión en terminal.
    """
    print(t("runtime.deliverable_id", value=submission.get("id", t("fallbacks.no_id"))))
    print(t("runtime.user_id", value=submission.get("userId", t("fallbacks.no_user_id"))))
    print(t("runtime.state", value=submission.get("state", t("fallbacks.no_state"))))
    print(t("runtime.late", value=submission.get("late", False)))
    print(t("runtime.assigned_grade", value=submission.get("assignedGrade")))
    print(t("runtime.draft_grade", value=submission.get("draftGrade")))
    print(t("runtime.resubmitted", value=is_resubmitted(submission)))
    print(t("runtime.ungraded", value=is_ungraded(submission)))
    print(t("runtime.attached", value=has_attachments(submission)))


def download_submission_attachments(
    submission: dict[str, Any],
    drive_service,
    submission_folder: str,
) -> list[str]:
    """
    Descarga attachments de una entrega.
    Regresa la items de archivos reales descargados.
    """
    attachments = get_attachments(submission)
    downloaded_paths: list[str] = []

    if not attachments:
        print(t("runtime.attachments_none"))
        return downloaded_paths

    print(t("runtime.attachments_header"))

    for att in attachments:
        if "driveFile" in att:
            drive_file = att.get("driveFile", {})
            drive_meta = drive_file.get("driveFile") or drive_file

            file_id = drive_meta.get("id")
            title = drive_meta.get("title", "archivo")
            mime_type = drive_meta.get("mimeType", "")

            print(f"    - DriveFile: {title} | id={file_id}")

            if file_id:
                path = download_file(
                    drive_service=drive_service,
                    file_id=file_id,
                    file_name=title,
                    folder=submission_folder,
                    mime_type=mime_type,
                )
                if path:
                    downloaded_paths.append(path)

        elif "link" in att:
            link = att["link"]
            print(t("runtime.link", title=link.get("title", t("fallbacks.no_title")), url=link.get("url", t("fallbacks.no_url"))))

        elif "form" in att:
            form = att["form"]
            print(t("runtime.form", title=form.get("title", t("fallbacks.no_title")), url=form.get("formUrl", t("fallbacks.no_url"))))

        elif "youTubeVideo" in att:
            video = att["youTubeVideo"]
            print(t("runtime.youtube", title=video.get("title", t("fallbacks.no_title")), url=video.get("alternateLink", t("fallbacks.no_url"))))

        else:
            print(t("runtime.unhandled_attachment"))

    return downloaded_paths


# ==========================================================
# CSV y ZIP
# ==========================================================

def write_summary_csv(csv_path: str, filas: list[dict[str, str]]) -> None:
    """
    Genera CSV general con el schema de salida final.
    """
    ensure_directory(os.path.dirname(csv_path))

    with open(csv_path, "w", newline="", encoding="utf-8") as csvfile:
        writer = csv.DictWriter(
            csvfile,
            fieldnames=CSV_OUTPUT_COLUMNS,
            extrasaction="ignore",
        )
        writer.writeheader()
        for fila in filas:
            row = {col: fila.get(col, "") for col in CSV_OUTPUT_COLUMNS}
            writer.writerow(row)

    print(f"\n{t('runtime.csv_generated', path=csv_path)}")


def compress_folder_to_zip(carpeta_origen: str, zip_sin_extension: str) -> str:
    """
    Comprime toda la carpeta en un zip.
    """
    zip_path = shutil.make_archive(zip_sin_extension, "zip", carpeta_origen)
    print(t("runtime.zip_generated", path=zip_path))
    return zip_path


# ==========================================================
# Procesamiento de una actividad
# ==========================================================

def process_activity(
    classroom_service,
    drive_service,
    course: dict[str, Any],
    coursework: dict[str, Any],
    download_mode: str,
    base_folder: str,
    profile_cache: dict[str, dict[str, str]],
    csv_rows: list[dict[str, str]],
    stats: dict[str, int],
    profile_scope_available: bool,
) -> None:
    """
    Procesa una actividad completa:
    - recupera entregas
    - aplica filter_name
    - descarga attachments
    - agrega filas al CSV
    - actualiza estadísticas
    """
    course_id = course["id"]
    course_name = course.get("name", f"{t('fallbacks.course')}_{course_id}")

    coursework_id = coursework["id"]
    coursework_title = coursework.get("title", f"{t('fallbacks.activity')}_{coursework_id}")
    coursework_display = coursework.get("display_name", coursework_title)

    print("\n" + "=" * 90)
    print(t("runtime.processing_activity", name=coursework_display))
    print("=" * 90)

    submissions = get_all_submissions(
        classroom_service=classroom_service,
        course_id=course_id,
        coursework_id=coursework_id,
    )

    stats["actividades_procesadas"] += 1
    stats["entregas_totales"] += len(submissions)

    if not submissions:
        print(t("runtime.no_submissions"))
        return

    entregas_filtradas = filter_submissions(submissions, download_mode)
    stats["entregas_filtradas"] += len(entregas_filtradas)

    print(t("runtime.submissions_total", count=len(submissions)))
    print(t("runtime.submissions_filtered", count=len(entregas_filtradas)))

    if not entregas_filtradas:
        print(t("runtime.submissions_no_match"))
        return

    nombre_carpeta_actividad = build_activity_slug(coursework_title, coursework_id)

    # Si la base_folder ya apunta exactamente a esta actividad
    # (caso descarga de una sola actividad), reutilízala tal cual.
    if os.path.basename(os.path.normpath(base_folder)) == nombre_carpeta_actividad:
        activity_folder = base_folder
    else:
        activity_folder = os.path.join(base_folder, nombre_carpeta_actividad)

    # Evita folders duplicados cuando se reprocesa la misma actividad.
    # Si la carpeta ya existe, se reconstruye limpia con la convención actual.
    if os.path.exists(activity_folder):
        shutil.rmtree(activity_folder)

    ensure_directory(activity_folder)

    for submission in entregas_filtradas:
        print_submission_summary(submission)

        user_id = submission.get("userId", t("fallbacks.no_user_id"))

        if user_id not in profile_cache:
            profile = {
                "correo": "",
                "name": "",
                "last_name": "",
            }

            if profile_scope_available:
                try:
                    profile = get_user_profile(
                        classroom_service=classroom_service,
                        user_id=user_id,
                        profile_scope_available=True,
                    )
                except HttpError as err:
                    print(
                        t("runtime.user_profile_fallback", user_id=user_id, err=err)
                    )
                    profile = extract_user_data_from_history(submission)

            else:
                profile = extract_user_data_from_history(submission)

            profile_cache[user_id] = profile

        profile = profile_cache[user_id]

        nombre_carpeta_entrega = build_submission_folder_name(
            submission=submission,
            profile=profile,
        )

        submission_folder = os.path.join(activity_folder, nombre_carpeta_entrega)
        ensure_directory(submission_folder)

        downloaded_paths: list[str] = []

        if can_download_submission(submission):
            downloaded_paths = download_submission_attachments(
                submission=submission,
                drive_service=drive_service,
                submission_folder=submission_folder,
            )

            if downloaded_paths:
                stats["archivos_descargados"] += len(downloaded_paths)

            evaluation = evaluate_submission_automatically(
                submission=submission,
                downloaded_paths=downloaded_paths,
                coursework=coursework,
                autograding_config=AUTOGRADING_CONFIG,
            )
        else:
            print(t("runtime.no_submission_download"))
            evaluation = {
                "auto_grade": 0,
                "auto_score": 0,
                "feedback": t("feedback.no_submission"),
                "auto_feedback": t("feedback.no_submission"),
                "auto_grading_reason": t("feedback.reason_no_submission"),
                "confidence_score": 1.0,
                "submission_type": labels()["submission_type"]["none"],
                "primary_file_type": "",
                "late_penalty": 0,
                "days_late": 0,
                "has_attachment": "false",
                "manual_review": "false",
                "requires_manual_review": "false",
                "readable_content": "false",
                "is_readable": "false",
                "delivery_valid_score": 0,
                "evidence_score": 0,
                "readable_content_score": 0,
                "minimum_sufficiency_score": 0,
                "keyword_hits": "",
                "keyword_hits_list": [],
                "content_score": 0,
                "files_read_for_content": 0,
                "detected_words": 0,
            }

        print(f"  auto_grade: {evaluation['auto_grade']}")
        print(f"  content_score: {evaluation['content_score']}")
        print(f"  readable_content: {evaluation['readable_content']}")
        print(f"  manual_review: {evaluation['manual_review']}")
        print(f"  late_penalty: {evaluation['late_penalty']}")

        csv_rows.append(
            {
                "course_name": course_name,
                "activity_name": coursework_title,
                "student_name": " ".join(
                    p for p in [profile.get("last_name", ""), profile.get("name", "")]
                    if p
                ).strip(),
                "student_mail": profile.get("correo", ""),
                "submission_status": get_readable_submission_status(submission),
                "has_attachment": str(evaluation["has_attachment"]).lower(),
                "submission_type": evaluation["submission_type"],
                "primary_file_type": evaluation.get("primary_file_type", ""),
                "days_late": str(evaluation["days_late"]),
                "is_readable": str(evaluation["is_readable"]).lower(),
                "word_count": str(evaluation["detected_words"]),
                "keyword_hits": evaluation["keyword_hits"],
                "requires_manual_review": str(evaluation["requires_manual_review"]).lower(),
                "confidence_score": f"{float(evaluation['confidence_score']):.2f}",
                "auto_score": str(evaluation["auto_score"]),
                "auto_feedback": normalize_basic_ascii(evaluation["auto_feedback"]),
                "auto_grading_reason": normalize_basic_ascii(evaluation["auto_grading_reason"]),
                "ai_feedback": "",
                "final_grade": "",
                "final_feedback": "",
            }
        )

        print("-" * 70)


# ==========================================================
# Flujo principal
# ==========================================================

AUTOGRADING_CONFIG = load_autograding_config()


# ==========================================================
# Build rubric runtime json
# ==========================================================

build_rubric_runtime_json(
    rubric_xlsx_path="config/Rubric.xlsx",
    output_json_path="config/rubric_runtime.json",
)


def main() -> None:
    """
    Flujo principal:
    1. autentica
    2. elige curso
    3. elige alcance
    4. filtra actividades
    5. elige filter_name de entregas
    6. elige formato de salida
    7. descarga
    8. genera CSV y zip

    """
    settings = get_settings()
    ensure_directories(settings)

    creds = get_credentials(
        credentials_path=settings.credentials_path,
        token_path=settings.token_path,
    )

    print(t("runtime.auth_ok"))
    print(t("runtime.token_valid", value=creds.valid))

    try:
        classroom_service = build("classroom", "v1", credentials=creds)
        drive_service = build("drive", "v3", credentials=creds)

        profile_scope_available = detect_profile_scope(classroom_service)
        if profile_scope_available:
            print(t("runtime.profile_scope_ok"))
        else:
            print(t("runtime.profile_scope_missing"))

        courses = get_all_courses(classroom_service)

        if not courses:
            print(t("runtime.no_courses"))
            return

        while True:
            selected_course = select_option(
                courses,
                t("menu.course"),
                allow_back=False,
            )
            if selected_course is None:
                continue

            course_id = selected_course["id"]
            course_slug = build_course_slug(
                selected_course.get("name", f"curso_{course_id}"),
                course_id,
            )
            course_display = selected_course.get(
                "display_name",
                selected_course.get("name", t("fallbacks.unnamed_course")),
            )

            print(f"\n{t('runtime.selected_course', value=course_display)}")

            while True:
                download_scope = select_download_scope(allow_back=True)
                if download_scope is None:
                    break

                print(f"\n{t('runtime.selected_scope', value=download_scope)}")

                while True:
                    activity_filter = select_activity_filter(allow_back=True)
                    if activity_filter is None:
                        break

                    print(f"\n{t('runtime.selected_filter', value=activity_filter)}")

                    while True:
                        download_mode = select_download_mode(allow_back=True)
                        if download_mode is None:
                            break

                        print(f"\n{t('runtime.selected_mode', value=describe_download_mode(download_mode))}")

                        while True:
                            output_format = select_output_format(allow_back=True)
                            if output_format is None:
                                break

                            print(t("runtime.selected_output", value=output_format))

                            courseworks = get_all_activities(classroom_service, course_id)

                            if not courseworks:
                                print(t("runtime.no_courseworks"))
                                break

                            filtered_courseworks = filter_activities(
                                classroom_service=classroom_service,
                                course_id=course_id,
                                courseworks=courseworks,
                                filter_name=activity_filter,
                            )

                            if not filtered_courseworks:
                                print(t("runtime.no_courseworks_after_filter"))
                                break

                            print(t("runtime.courseworks_found", total=len(courseworks), filtered=len(filtered_courseworks)))

                            profile_cache: dict[str, dict[str, str]] = {}
                            csv_rows: list[dict[str, str]] = []
                            stats = {
                                "actividades_procesadas": 0,
                                "entregas_totales": 0,
                                "entregas_filtradas": 0,
                                "archivos_descargados": 0,
                            }

                            export_directory = get_export_directory(settings)
                            course_folder = os.path.normpath(
                                os.path.join(export_directory, course_slug)
                            )

                            if download_scope == "single_coursework":
                                selected_coursework = select_option(
                                    filtered_courseworks,
                                    t("menu.activity"),
                                    allow_back=True,
                                )
                                if selected_coursework is None:
                                    continue

                                print(f"\n{t('runtime.selected_activity', value=selected_coursework['display_name'])}")

                                confirmado = confirm_download_summary(
                                    course_display=course_display,
                                    download_scope=download_scope,
                                    activity_filter=activity_filter,
                                    download_mode=download_mode,
                                    output_format=output_format,
                                    total_actividades=1,
                                    selected_activity=selected_coursework,
                                )
                                if not confirmado:
                                    continue

                                base_folder = prepare_output_directory(
                                    course_folder,
                                    limpiar_si_existe=True,
                                )

                                process_activity(
                                    classroom_service=classroom_service,
                                    drive_service=drive_service,
                                    course=selected_course,
                                    coursework=selected_coursework,
                                    download_mode=download_mode,
                                    base_folder=base_folder,
                                    profile_cache=profile_cache,
                                    csv_rows=csv_rows,
                                    stats=stats,
                                    profile_scope_available=profile_scope_available,
                                )

                                nombre_csv = "resumen_entregas.csv"
                                nombre_zip = course_slug

                            elif download_scope == "all_courseworks":
                                confirmado = confirm_download_summary(
                                    course_display=course_display,
                                    download_scope=download_scope,
                                    activity_filter=activity_filter,
                                    download_mode=download_mode,
                                    output_format=output_format,
                                    total_actividades=len(filtered_courseworks),
                                    selected_activity=None,
                                )
                                if not confirmado:
                                    continue

                                base_folder = prepare_output_directory(
                                    course_folder,
                                    limpiar_si_existe=True,
                                )

                                print(f"\n{t('runtime.processing_all_courseworks', count=len(filtered_courseworks))}")

                                for idx, coursework in enumerate(filtered_courseworks, start=1):
                                    print(f"\n[{idx}/{len(filtered_courseworks)}] {coursework['display_name']}")

                                    process_activity(
                                        classroom_service=classroom_service,
                                        drive_service=drive_service,
                                        course=selected_course,
                                        coursework=coursework,
                                        download_mode=download_mode,
                                        base_folder=base_folder,
                                        profile_cache=profile_cache,
                                        csv_rows=csv_rows,
                                        stats=stats,
                                        profile_scope_available=profile_scope_available,
                                    )

                                nombre_csv = "resumen_todas_las_actividades.csv"
                                nombre_zip = course_slug

                            else:
                                print(t("runtime.unknown_scope"))
                                continue

                            csv_path = os.path.join(base_folder, nombre_csv)
                            write_summary_csv(csv_path, csv_rows)

                            if output_format == "zip_and_folder":
                                zip_base_name = os.path.join(export_directory, nombre_zip)
                                compress_folder_to_zip(base_folder, zip_base_name)

                            print("\n" + "=" * 90)
                            print(t("runtime.final_summary"))
                            print("=" * 90)
                            print(f"{t('ui.course')}: {course_display}")
                            print(t("runtime.activities_processed", count=stats["actividades_procesadas"]))
                            print(t("runtime.submissions_seen", count=stats["entregas_totales"]))
                            print(t("runtime.submissions_filter_matched", count=stats["entregas_filtradas"]))
                            print(t("runtime.files_downloaded", count=stats["archivos_descargados"]))
                            print(t("runtime.csv_rows", count=len(csv_rows)))
                            print(t("runtime.all_mode_note"))
                            print(t("runtime.base_folder", path=base_folder))
                            if output_format == "zip_and_folder":
                                print(t("runtime.zip_path", path=f"{os.path.join(export_directory, nombre_zip)}.zip"))

                            print(f"\n{t('runtime.process_done')}")
                            return

    except HttpError as err:
        print(t("runtime.classroom_error", err=err))


if __name__ == "__main__":
    main()
